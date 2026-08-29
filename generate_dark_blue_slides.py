import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen layout
prs.slide_height = Inches(7.5)

# Styling Palette Constants (Premium Widescreen Dark Navy/Cyan Theme)
COLOR_BG = RGBColor(11, 22, 38)            # Deep dark navy base background (#0b1626)
COLOR_CARD_BG = RGBColor(22, 36, 56)        # Slightly lighter dark blue for cards (#162438)
COLOR_TEXT_MAIN = RGBColor(255, 255, 255)   # Pure White main text (#ffffff)
COLOR_TEXT_SUB = RGBColor(190, 200, 218)    # Light gray-blue secondary text (#bec8da)
COLOR_ACCENT_CYAN = RGBColor(0, 180, 216)   # Neon Cyan/Teal for highlights (#00b4d8)
COLOR_BORDER = RGBColor(38, 56, 82)         # Muted dark blue for borders (#263852)

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header(slide, title_text, category_text="BMSCE CAMPUS SOCIAL PLATFORM"):
    # Add top-left "Team:commitnpray" indicator
    txTeam = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(4.0), Inches(0.4))
    tfTeam = txTeam.text_frame
    pTeam = tfTeam.paragraphs[0]
    pTeam.text = "Team:commitnpray"
    pTeam.font.name = "Arial"
    pTeam.font.size = Pt(11)
    pTeam.font.bold = True
    pTeam.font.color.rgb = COLOR_ACCENT_CYAN
    pTeam.font.underline = True

    # Add category tag / section tracker
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.7), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = category_text.upper()
    p.font.name = "Arial"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_CYAN
    
    # Add main slide title
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = title_text
    p2.font.name = "Arial"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_MAIN

def create_card(slide, x, y, width, height, title_text, body_text, border_color=COLOR_BORDER, title_color=COLOR_ACCENT_CYAN):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CARD_BG
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.space_after = Pt(5)
    
    p2 = tf.add_paragraph()
    p2.text = body_text
    p2.font.name = "Arial"
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = COLOR_TEXT_MAIN
    
    return shape

def add_arrow(slide, x, y, width, height):
    # Renders a sleek neon cyan arrow to represent flow paths
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT_CYAN
    shape.line.color.rgb = COLOR_ACCENT_CYAN
    shape.line.width = Pt(1)
    return shape

def add_placeholder_box(slide, x, y, width, height, label_text):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(16, 29, 48)
    shape.line.color.rgb = COLOR_ACCENT_CYAN
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label_text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_SUB
    return shape

# ----------------- SLIDE 1: PITCH DECK TITLE SLIDE -----------------
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)

# Team name at top-left
txTeam = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(4.0), Inches(0.4))
tfTeam = txTeam.text_frame
pTeam = tfTeam.paragraphs[0]
pTeam.text = "Team:commitnpray"
pTeam.font.name = "Arial"
pTeam.font.size = Pt(12)
pTeam.font.bold = True
pTeam.font.color.rgb = COLOR_ACCENT_CYAN
pTeam.font.underline = True

# Main Center Title Box
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(3.0))
tf = txBox.text_frame
tf.word_wrap = True

p_title = tf.paragraphs[0]
p_title.text = "BMSCE Social & Intelligence Platform"
p_title.alignment = PP_ALIGN.CENTER
p_title.font.name = "Arial"
p_title.font.size = Pt(44)
p_title.font.bold = True
p_title.font.color.rgb = COLOR_TEXT_MAIN
p_title.space_after = Pt(14)

p_sub = tf.add_paragraph()
p_sub.text = "Connecting students across the campus through AI-driven discovery."
p_sub.alignment = PP_ALIGN.CENTER
p_sub.font.name = "Arial"
p_sub.font.size = Pt(18)
p_sub.font.color.rgb = COLOR_TEXT_SUB

# Bottom Left Author Box
txAuthor = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(4.0), Inches(1.5))
tfAuthor = txAuthor.text_frame
tfAuthor.word_wrap = True

p_name = tfAuthor.paragraphs[0]
p_name.text = "Kritika Panwar"
p_name.font.name = "Arial"
p_name.font.size = Pt(15)
p_name.font.bold = True
p_name.font.color.rgb = COLOR_TEXT_MAIN
p_name.space_after = Pt(4)

p_reg = tfAuthor.add_paragraph()
p_reg.text = "1BF24CS151"
p_reg.font.name = "Arial"
p_reg.font.size = Pt(13)
p_reg.font.color.rgb = COLOR_TEXT_SUB

# ----------------- SLIDE 2: THE PROBLEM: A DISCONNECTED CAMPUS -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "The Problem: A Disconnected Campus", "01. PROBLEM DEFINITION")

# Question banner in Cyan
txBanner = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.8))
tfBanner = txBanner.text_frame
tfBanner.word_wrap = True
pBanner = tfBanner.paragraphs[0]
pBanner.text = "How can we create a single intelligent campus platform that connects students with the places, people and information they need to make the most of campus life?"
pBanner.font.name = "Arial"
pBanner.font.size = Pt(14.5)
pBanner.font.bold = True
pBanner.font.color.rgb = COLOR_ACCENT_CYAN

# Left Column Bullets
txBullets = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(6.2), Inches(4.5))
tfBullets = txBullets.text_frame
tfBullets.word_wrap = True

bullets = [
    ("Large & Diverse Campus", "Nearly 15 acres with 18 undergraduate programs, making campus navigation and discovery challenging."),
    ("Difficult Navigation", "Freshers often struggle to find unfamiliar blocks, floors, labs and facilities across campus."),
    ("Limited Cross-Department Connections", "Students have few ways to discover peers with shared interests, skills and goals beyond their own departments."),
    ("Scattered Information", "Campus information, student experiences and opportunities are spread across different sources.")
]

for idx, (title, text) in enumerate(bullets):
    p = tfBullets.add_paragraph() if idx > 0 else tfBullets.paragraphs[0]
    p.text = f"•  {title}: "
    p.font.name = "Arial"
    p.font.size = Pt(13.5)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    
    # Body inline
    p.text += text
    p.space_after = Pt(16)

# Right Column Side Cards
create_card(slide, 7.3, 2.6, 5.2, 2.0, 
            "📍 01 — FINDING YOUR WAY", 
            "\"Where is that lab?\"\nFor freshers, finding unfamiliar blocks, floors, labs and facilities across campus can be confusing.",
            title_color=COLOR_ACCENT_CYAN)

create_card(slide, 7.3, 4.8, 5.2, 2.0, 
            "👥 02 — FINDING YOUR PEOPLE", 
            "\"Who can I build this with?\"\nStudents often struggle to discover peers beyond their department for hackathons, projects, clubs and campus events.",
            title_color=COLOR_ACCENT_CYAN)

# ----------------- SLIDE 3: SCOPE: WHAT WE ARE BUILDING -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Scope: What We Are Building", "02. THE SOLUTION")

scopes = [
    ("Gamified Campus Exploration Scanner", "A mobile-first camera scanning portal letting users scan landmarks, verify locations, and unlock coordinate cards under distinct Rarity Tiers."),
    ("Fog-of-War Interactive Campus Map", "An interactive map start layout that begins greyed out and unlocks sector-by-sector as students photograph blocks, making campus navigation a game."),
    ("Campus Passport Stamp Book", "A collection log that stamps the user's book with official inked seals when they visit campus zones, unlocking badges for set completions."),
    ("Interest-Based Matchmaker", "Social matching displaying overlap metrics (e.g. 85% Match) for peer skills, prompting chat requests and coordinate pin sharing."),
    ("Databricks Genie AI Insights Desk", "An analytics desk allowing administrators to ask plain text questions and dynamically generate traffic heatmaps of the campus.")
]

y_start = 1.8
for idx, (title, desc) in enumerate(scopes):
    create_card(slide, 0.8, y_start, 11.7, 0.9, 
                title, desc, border_color=COLOR_BORDER, title_color=COLOR_ACCENT_CYAN)
    y_start += 1.05

# ----------------- SLIDE 4: BMSCE PLATFORM ECOSYSTEM -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "BMSCE Platform Ecosystem", "03. FUNCTIONAL FLOW")

create_card(slide, 0.8, 1.8, 5.6, 1.5, 
            "Mobile Front-End Web Client", 
            "A fast React view wrapper with five central dashboards: Explorer (scanning/fog map), Campus (matching), Genie (chatbot), Inbox (messages), and Profile.",
            title_color=COLOR_ACCENT_CYAN)

create_card(slide, 0.8, 3.5, 5.6, 1.5, 
            "Serverless Backend Controller", 
            "Node.js API middleware handles file compression, geofence verification, explore streaks comparison, and messaging gateways.",
            title_color=COLOR_ACCENT_CYAN)

create_card(slide, 0.8, 5.2, 5.6, 1.5, 
            "Lakehouse Analytics Sync", 
            "Google Firestore logs are periodically synced to a Delta Lakehouse model, enabling admins to generate charts using Databricks Genie.",
            title_color=COLOR_ACCENT_CYAN)

add_placeholder_box(slide, 6.9, 1.8, 5.6, 4.9, "[ ECOSYSTEM INTERFACES ARCHITECTURE PREVIEW ]")

# ----------------- SLIDE 5: USER JOURNEY: EXPLORING THE CAMPUS -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "User Journey: Exploring the Campus", "04. GAMIFIED NAVIGATION")

steps = [
    ("1. Navigate Gray Zones", "Student opens Fog Map and sees the Engineering Block C coordinates locked in gray fog."),
    ("2. Snap Landmark Photo", "Student visits the physical block, opens the Scanner, and photographs the Robotics Lab."),
    ("3. Stamp Virtual Passport", "Firestore logs the stamp, stamps the Passport with an ink seal, and unlocks the map tile."),
    ("4. Promote Level & Rank", "Scan grants +120 XP (First to Document) triggering the level promotion overlay.")
]

y_step = 1.8
for idx, (title, desc) in enumerate(steps):
    create_card(slide, 0.8, y_step, 5.6, 1.15, title, desc, title_color=COLOR_ACCENT_CYAN)
    y_step += 1.25

add_placeholder_box(slide, 6.9, 1.8, 5.6, 4.9, "[ SCREENSHOT: EXPLORER SCANNER / FOG MAP / PASSPORT ]")

# ----------------- SLIDE 6: USER JOURNEY: FINDING YOUR TRIBE -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "User Journey: Finding Your Tribe", "05. SOCIAL MATCHMAKING")

social_steps = [
    ("1. Enter Interest Tags", "Student sets profile avatar and chooses interest tags (AI, React, Figma)."),
    ("2. Verify Match Overlaps", "Campus matches sort classmates based on skill overlap percentages (e.g. 85% Match)."),
    ("3. Start Direct Messages", "Student connects with a peer, initiating chat channels and image sharing attachments."),
    ("4. Share Landmark Pins", "Student shares the verified coordinate pin of a study spot, letting the peer view it instantly.")
]

y_step = 1.8
for idx, (title, desc) in enumerate(social_steps):
    create_card(slide, 0.8, y_step, 5.6, 1.15, title, desc, title_color=COLOR_ACCENT_CYAN)
    y_step += 1.25

add_placeholder_box(slide, 6.9, 1.8, 5.6, 4.9, "[ SCREENSHOT: PEER MATCHES / DIRECT MESSAGES / INBOX ]")

# ----------------- SLIDE 7: ARCHITECTURE DIAGRAM (FULLY DRAWN FLOW) -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "System Architecture Blueprint", "06. TECHNICAL TOPOLOGY")

# Replaced text cards + placeholders with a full 5-stage horizontal flowchart mapping the entire stack
box_y = 2.8
box_height = 2.4
box_width = 1.85
arrow_y = 3.75
arrow_width = 0.45
arrow_height = 0.5

# 1. Mobile Client
create_card(slide, 0.8, box_y, box_width, box_height, 
            "📱 Client Tier", 
            "Next.js 14 / React\n• Responsive layouts\n• Caches views locally\n• Offline-first queueing\n• Hosted on Vercel")

add_arrow(slide, 2.65, arrow_y, arrow_width, arrow_height)

# 2. Server API
create_card(slide, 3.1, box_y, box_width, box_height, 
            "⚙️ Express API", 
            "Node.js Middleware\n• Validates coordinates\n• Checks checkin streaks\n• Routes direct messages\n• Hosted on Render")

add_arrow(slide, 4.95, arrow_y, arrow_width, arrow_height)

# 3. Firestore & Storage
create_card(slide, 5.4, box_y, box_width, box_height, 
            "🔥 Cloud Storage", 
            "Google Firestore\n• JSON Document DBs\n• Firebase Auth checks\n• Cloudinary photo hosts\n• Serverless scaling")

add_arrow(slide, 7.25, arrow_y, arrow_width, arrow_height)

# 4. GCS Delta Data Lake
create_card(slide, 7.7, box_y, box_width, box_height, 
            "🌊 Data Lake", 
            "GCS / Delta Lake\n• Sync extension exports\n• Stores daily updates\n• Mounted as Delta tables\n• Google Storage host")

add_arrow(slide, 9.55, arrow_y, arrow_width, arrow_height)

# 5. Databricks Genie
create_card(slide, 10.0, box_y, 2.5, box_height, 
            "🧠 Analytics", 
            "Databricks Genie CE\n• Natural language inputs\n• Auto-SQL translation\n• Renders traffic maps\n• Real-time admin views")

# ----------------- SLIDE 8: DATA/USER FLOW DIAGRAM (FULLY DRAWN FLOW) -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Data & User Flow Diagram", "07. PIPELINE FLOW")

# Replaced placeholder with a full 4-stage horizontal flow diagram mapping the location upload pipeline
flow_y = 2.8
flow_height = 2.4
flow_width = 2.2
arrow_flow_y = 3.75
arrow_flow_width = 0.55
arrow_flow_height = 0.5

# Step 1
create_card(slide, 0.8, flow_y, flow_width, flow_height, 
            "📸 1. Capture & Scan", 
            "Student takes a campus landmark photo. The Next.js client attaches GPS metadata and dispatches it securely to the Node server.",
            title_color=COLOR_ACCENT_CYAN)

add_arrow(slide, 3.0, arrow_flow_y, arrow_flow_width, arrow_flow_height)

# Step 2
create_card(slide, 3.55, flow_y, flow_width, flow_height, 
            "📦 2. Host & Save", 
            "Node server uploads the photo to Cloudinary. Once complete, it saves the optimized HTTPS image URL to Google Firestore.",
            title_color=COLOR_ACCENT_CYAN)

add_arrow(slide, 5.75, arrow_flow_y, arrow_flow_width, arrow_flow_height)

# Step 3
create_card(slide, 6.3, flow_y, flow_width, flow_height, 
            "🔥 3. Verify Checkins", 
            "Firestore checks location data, stamps the student's Passport document, increments streaks, and unlocks map tiles.",
            title_color=COLOR_ACCENT_CYAN)

add_arrow(slide, 8.5, arrow_flow_y, arrow_flow_width, arrow_flow_height)

# Step 4
create_card(slide, 9.05, 3.5, flow_y, flow_height, 
            "🏆 4. Award XP", 
            "API calculates points (+120 XP for First-ever scan) and unlocks achievements, showing a Level-Up overlay on the phone screen.",
            title_color=COLOR_ACCENT_CYAN)

# ----------------- SLIDE 9: DATA MODEL USAGE -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Data Model Usage", "08. DATABASE DESIGN")

col_width = 3.65
x_col1 = 0.8
x_col2 = 4.85
x_col3 = 8.9
height_col = 4.9

create_card(slide, x_col1, 1.8, col_width, height_col, 
            "Student Collections", 
            "•  users collection\nPath: `/users/{userId}`\nFields: name, username, avatarUrl, level, xp, interests, skills, exploreStreak, lastScanDate\n\n•  passport_stamps collection\nPath: `/users/{userId}/stamps/{stampId}`\nFields: locationId, unlockedAt, isFirstDocumenter, locationName",
            title_color=COLOR_TEXT_MAIN)

create_card(slide, x_col2, 1.8, col_width, height_col, 
            "Campus Collections", 
            "•  locations collection\nPath: `/locations/{locationId}`\nFields: name, building, floor, coordinates, rarity, facilities, tips\n\n•  map_tiles collection\nPath: `/map_tiles/{tileId}`\nFields: name, coordinates, iconType, associatedLocationId",
            title_color=COLOR_TEXT_MAIN)

create_card(slide, x_col3, 1.8, col_width, height_col, 
            "Social & Community", 
            "•  messages collection\nPath: `/chats/{chatId}/messages/{messageId}`\nFields: senderId, receiverId, text, imageUrl, sharedLocationId, timestamp\n\n•  posts collection\nPath: `/posts/{postId}`\nFields: authorName, title, content, tag, likes, timestamp",
            title_color=COLOR_TEXT_MAIN)

# ----------------- SLIDE 10: AI FEATURES AND GROUNDING -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "AI Features & Grounding Model", "09. ARTIFICIAL INTELLIGENCE")

create_card(slide, x_col1, 1.8, col_width, height_col, 
            "Databricks Genie Agent", 
            "Provides administrators with conversational SQL data analysis. Genie maps natural English to mounted database schemas, automatically translating queries like 'Which landmarks had the most scans?' into charts without coding queries.",
            title_color=COLOR_ACCENT_CYAN)

create_card(slide, x_col2, 1.8, col_width, height_col, 
            "Gemini Grounded Chatbot", 
            "Powers the Genie Chatbot in the mobile client. It is grounded by system instructions seeded with campus guides, print rates, office locations, and schedules, answering student requests accurately.",
            title_color=COLOR_ACCENT_CYAN)

create_card(slide, x_col3, 1.8, col_width, height_col, 
            "Profile Matchmaker", 
            "Uses mathematical sets calculations directly on Firestore collections. Computes overlap ratios across classmate interests, skills, and hobbies to matches students with high-score partners.",
            title_color=COLOR_ACCENT_CYAN)

# ----------------- SLIDE 11: VALUE PROPOSITION AND BENEFITS -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Value Proposition & Benefits", "10. STRATEGIC VALUE")

create_card(slide, x_col1, 1.8, col_width, height_col, 
            "For Students", 
            "•  Onboarding Speed: Freshmen discover campus facilities (printer spots, hardware labs) in minutes.\n\n•  Peer Matchmaking: Connects peers directly based on code skills or hobbies for hackathon sprints.\n\n•  Active Play: Gamified achievements make daily walks fun.",
            title_color=COLOR_TEXT_MAIN)

create_card(slide, x_col2, 1.8, col_width, height_col, 
            "For Admins", 
            "•  Natural Queries: Admins generate charts dynamically using Databricks Genie.\n\n•  Analytics: Reveals real-time spatial heatmaps showing popular spots.\n\n•  Resource Planning: Helps coordinates schedules based on facility traffic checks.",
            title_color=COLOR_TEXT_MAIN)

create_card(slide, x_col3, 1.8, col_width, height_col, 
            "For the Institution", 
            "•  Student Retention: Unifies student communication, reducing early isolation.\n\n•  Platform Consolidation: Integrates geolocations, chats, and forums into one app.\n\n•  Event Engagement: Drives participation in scavenger hunts and tech club events.",
            title_color=COLOR_TEXT_MAIN)

# ----------------- SLIDE 12: TECH STACK & PRODUCTION PLATFORMS -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Technology Stack & Free Tier Plan", "11. CLOUD STACK")

tech_cards = [
    ("Next.js & React (Vercel Host)", "Deploys static assets. Hobby free tier gives 100 GB monthly bandwidth and free SSL domain connection."),
    ("Google Firestore (Firebase Free Tier)", "Serverless NoSQL database. Free tier includes 50,000 reads, 20,000 writes daily, and 1 GiB storage."),
    ("Node.js & Express (Render Host)", "Runs backend geofence controls, explore streaks calculations, and real-time message streams."),
    ("Cloudinary (Media Hosting)", "Compresses and stores student scans. Free tier includes 25,000 monthly image transformations or 25 GB storage."),
    ("Databricks Genie Agent (CE)", "Databricks Community Edition processes analytics tables mounted from Firestore data lakes.")
]

y_tech = 1.8
for idx, (title, desc) in enumerate(tech_cards):
    create_card(slide, 0.8, y_tech, 11.7, 0.9, title, desc, title_color=COLOR_ACCENT_CYAN)
    y_tech += 1.05

# ----------------- SLIDE 13: EXECUTION PLAN - 12 HOURS -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Execution Plan - 12 Hours", "12. ROADMAP")

timeline = [
    ("H 1 - 3: Setup & Database", "Config Next.js app, style theme files, set up Firebase Firestore document collections, and initialize auth APIs."),
    ("H 4 - 6: Backend & Upload", "Code backend Express APIs for photo scans, geofencing validations, Cloudinary storage uploads, and streak comparators."),
    ("H 7 - 9: Client Integration", "Bind React Context state logic inside views. Integrate Maps coordinates, messaging streams, and Hub profile matches."),
    ("H 10 - 12: Analytics & Host", "Sync Firestore to GCS, mount Delta tables in Databricks Genie, run compile builds, and deploy onto Vercel and Render.")
]

x_time = 0.8
width_time = 2.75
gap = 0.2
for idx, (title, desc) in enumerate(timeline):
    create_card(slide, x_time, 1.8, width_time, 4.9, title, desc, title_color=COLOR_ACCENT_CYAN)
    x_time += width_time + gap

# ----------------- SLIDE 14: RISKS AND FALLBACKS -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Risks and Fallbacks", "13. MITIGATION")

create_card(slide, x_col1, 1.8, col_width, height_col, 
            "Risk 1: Firestore Daily Read limits", 
            "•  Impact: App crashes once reads surpass the daily 50k limit under high student usage.\n\n•  Mitigation: Implement aggressive client-side caching via browser LocalStorage. Only fetch database collections when new chat pins or scans are updated, keeping reads minimal.",
            border_color=COLOR_BORDER, title_color=COLOR_TEXT_MAIN)

create_card(slide, x_col2, 1.8, col_width, height_col, 
            "Risk 2: Render Free Tier Cold Starts", 
            "•  Impact: Rendering scanner triggers lags 30-50 seconds when user server restarts.\n\n•  Mitigation: Send a periodic uptime heartbeat query from the client application background to keep Render's free instance active during peak campus hours.",
            border_color=COLOR_BORDER, title_color=COLOR_TEXT_MAIN)

create_card(slide, x_col3, 1.8, col_width, height_col, 
            "Risk 3: Internet Outage in Basements", 
            "•  Impact: Camera scans fail in basement labs due to dead cellular signals.\n\n•  Mitigation: Buffer scans in a local offline queue in LocalStorage, syncing coordinates and images once internet connections are restored.",
            border_color=COLOR_BORDER, title_color=COLOR_TEXT_MAIN)

# Save presentation
prs.save("BMSCE_Campus_Social_Platform_DarkBlue.pptx")
print("Fully drawn Presentation generated successfully!")
