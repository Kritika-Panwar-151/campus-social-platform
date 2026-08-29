import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen layout
prs.slide_height = Inches(7.5)

# Styling Palette Constants (Premium Gamified Light Theme)
COLOR_BG = RGBColor(248, 247, 255)           # Ermine White base background (#f8f7ff)
COLOR_CARD_BG = RGBColor(255, 255, 255)      # Pure White card backgrounds (#ffffff)
COLOR_TEXT_MAIN = RGBColor(51, 49, 54)       # Stiletto Charcoal dark grey (#333136)
COLOR_TEXT_SUB = RGBColor(99, 93, 115)       # Luscious Purple secondary text (#635d73)
COLOR_ACCENT_ORANGE = RGBColor(242, 143, 95)  # Candied Yam Tangerine (#f28f5f)
COLOR_ACCENT_PURPLE = RGBColor(162, 124, 248) # Purple Illusionist accent (#a27cf8)
COLOR_BORDER = RGBColor(197, 186, 232)        # Muted purple for card borders (#c5bae8)

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header(slide, title_text, category_text="BMSCE CAMPUS SOCIAL PLATFORM"):
    # Add category tag / section tracker
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = category_text.upper()
    p.font.name = "Arial"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_ORANGE
    
    # Add main slide title
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = title_text
    p2.font.name = "Arial"
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_MAIN

def create_card(slide, x, y, width, height, title_text, body_text, border_color=COLOR_BORDER, title_color=COLOR_ACCENT_PURPLE):
    # Draw a clean rounded rectangle to represent a premium UI card
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CARD_BG
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    
    # Add text inside the card
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    
    # Card Title
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Arial"
    p.font.size = Pt(13.5)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.space_after = Pt(6)
    
    # Card Body
    p2 = tf.add_paragraph()
    p2.text = body_text
    p2.font.name = "Arial"
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_TEXT_MAIN
    
    return shape

def add_placeholder_box(slide, x, y, width, height, label_text):
    # Screenshot / diagram outline placeholder
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 238, 250)
    shape.line.color.rgb = COLOR_ACCENT_PURPLE
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label_text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_SUB
    return shape

# ----------------- SLIDE 1: PITCH DECK TITLE SLIDE -----------------
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)

# Accent background card for title
title_card = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(1.5), Inches(11.73), Inches(4.5)
)
title_card.fill.solid()
title_card.fill.fore_color.rgb = COLOR_CARD_BG
title_card.line.color.rgb = COLOR_ACCENT_PURPLE
title_card.line.width = Pt(2.5)

tf_title = title_card.text_frame
tf_title.word_wrap = True
tf_title.margin_left = Inches(0.6)
tf_title.margin_top = Inches(0.8)

p_category = tf_title.paragraphs[0]
p_category.text = "CAMPUS INTELLECTUAL PLATFORM"
p_category.font.name = "Arial"
p_category.font.size = Pt(12)
p_category.font.bold = True
p_category.font.color.rgb = COLOR_ACCENT_ORANGE
p_category.space_after = Pt(12)

p_title = tf_title.add_paragraph()
p_title.text = "BMSCE Campus Social & Intelligence Platform"
p_title.font.name = "Arial"
p_title.font.size = Pt(40)
p_title.font.bold = True
p_title.font.color.rgb = COLOR_TEXT_MAIN
p_title.space_after = Pt(8)

p_sub = tf_title.add_paragraph()
p_sub.text = "A Gamified, AI-Powered Offline-First Campus Ecosystem"
p_sub.font.name = "Arial"
p_sub.font.size = Pt(16)
p_sub.font.color.rgb = COLOR_TEXT_SUB

# ----------------- SLIDE 2: THE PROBLEM: A DISCONNECTED CAMPUS (2x2 Card Grid) -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "The Problem: A Disconnected Campus", "01. PROBLEM DEFINITION")

# 2x2 grid coordinates
x_left = 0.8
x_right = 6.9
width = 5.6
height = 2.1
y_row1 = 1.8
y_row2 = 4.3

create_card(slide, x_left, y_row1, width, height, 
            "Freshman Isolation & Friction", 
            "Incoming students struggle to adapt to the geographic layout, classroom locations, and lab blocks. This navigation friction delays integration and makes early networking highly challenging.")

create_card(slide, x_right, y_row1, width, height, 
            "Fragmented Notice Channels", 
            "Announcements, tech workshops, hackathons, and printing facilities guidelines are scattered across emails, chat groups, and physical bulletin boards. Important events are frequently missed.")

create_card(slide, x_left, y_row2, width, height, 
            "Ad-Hoc Peer Matching", 
            "Forming teams for projects or hackathons is based on ad-hoc friendships rather than matching skills. Finding a React frontend dev or Python AI partner is highly unstructured.")

create_card(slide, x_right, y_row2, width, height, 
            "Underutilized Campus Facilities", 
            "Advanced resources like GPU computation clusters, campus libraries, and specialized student hubs remain underutilized because static PDFs do not encourage or guide exploration.")

# ----------------- SLIDE 3: SCOPE: WHAT WE ARE BUILDING (Structured List) -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Scope: What We Are Building", "02. THE SOLUTION")

scopes = [
    ("Gamified Coordinates Scanner", "A mobile-first web camera that lets users scan landmarks, verify locations, and unlock coordinate cards divided into distinct Rarity Tiers."),
    ("Interactive Fog-of-War Map", "A map start layout that begins greyed out and unlocks sector-by-sector as students photograph blocks, making campus navigation a game."),
    ("Campus Passport Stamp Book", "A collection log that stamps the user's book with official inked seals when they visit campus zones, unlocking badges for set completions."),
    ("Interest-Based Matchmaker", "Social matching displaying overlap metrics (e.g. 85% Match) for peer skills, prompting chat requests and coordinate pin sharing."),
    ("Databricks Genie AI Insights Desk", "An analytics desk allowing administrators to ask plain text questions and dynamically generate traffic heatmaps of the campus.")
]

y_start = 1.8
for idx, (title, desc) in enumerate(scopes):
    create_card(slide, 0.8, y_start, 11.7, 0.9, 
                title, desc, border_color=COLOR_BORDER, title_color=COLOR_ACCENT_ORANGE)
    y_start += 1.05

# ----------------- SLIDE 4: BMSCE PLATFORM ECOSYSTEM (Split Column) -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "BMSCE Platform Ecosystem", "03. FUNCTIONAL FLOW")

# Left Column (Ecosystem cards)
create_card(slide, 0.8, 1.8, 5.6, 1.5, 
            "Mobile Front-End Web Client", 
            "A fast React view wrapper with five central dashboards: Explorer (scanning/fog map), Campus (matching), Genie (chatbot), Inbox (messages), and Profile.",
            title_color=COLOR_ACCENT_ORANGE)

create_card(slide, 0.8, 3.5, 5.6, 1.5, 
            "Serverless Backend Controller", 
            "Node.js API middleware handles file compression, geofence verification, explore streaks comparison, and messaging gateways.",
            title_color=COLOR_ACCENT_ORANGE)

create_card(slide, 0.8, 5.2, 5.6, 1.5, 
            "Lakehouse Analytics Sync", 
            "Google Firestore logs are periodically synced to a Delta Lakehouse model, enabling admins to generate charts using Databricks Genie.",
            title_color=COLOR_ACCENT_ORANGE)

# Right Column (Placeholder for diagram)
add_placeholder_box(slide, 6.9, 1.8, 5.6, 4.9, "[ ECOSYSTEM INTERFACES ARCHITECTURE PREVIEW ]")

# ----------------- SLIDE 5: USER JOURNEY: EXPLORING THE CAMPUS (Step Flow) -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "User Journey: Exploring the Campus", "04. GAMIFIED NAVIGATION")

steps = [
    ("1. Navigate Gray Zones", "Student opens Fog Map and sees the Engineering Block C coordinates locked in gray fog."),
    ("2. Snap Landmark Photo", "Student visits the physical block, opens the Scanner, and photographs the Robotics Lab."),
    ("3. Stamp Virtual Passport", "Firestore logs the stamp, stamps the Passport with an ink seal, and unlocks the map tile."),
    ("4. Promote Level & Rank", "Scan grants +120 XP (First to Document) triggering the level promotion overlay.")
]

y_step = 1.8
for idx, (title, desc) in enumerate(steps):
    create_card(slide, 0.8, y_step, 5.6, 1.15, title, desc, title_color=COLOR_ACCENT_PURPLE)
    y_step += 1.25

add_placeholder_box(slide, 6.9, 1.8, 5.6, 4.9, "[ SCREENSHOT: EXPLORER SCANNER / FOG MAP / PASSPORT ]")

# ----------------- SLIDE 6: USER JOURNEY: FINDING YOUR TRIBE (Step Flow) -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "User Journey: Finding Your Tribe", "05. SOCIAL MATCHMAKING")

social_steps = [
    ("1. Enter Interest Tags", "Student sets profile avatar and chooses interest tags (AI, React, Figma)."),
    ("2. Verify Match Overlaps", "Campus matches sort classmates based on skill overlap percentages (e.g. 85% Match)."),
    ("3. Start Direct Messages", "Student connects with a peer, initiating chat channels and image sharing attachments."),
    ("4. Share Landmark Pins", "Student shares the verified coordinate pin of a study spot, letting the peer view it instantly.")
]

y_step = 1.8
for idx, (title, desc) in enumerate(social_steps):
    create_card(slide, 0.8, y_step, 5.6, 1.15, title, desc, title_color=COLOR_ACCENT_PURPLE)
    y_step += 1.25

add_placeholder_box(slide, 6.9, 1.8, 5.6, 4.9, "[ SCREENSHOT: PEER MATCHES / DIRECT MESSAGES / INBOX ]")

# ----------------- SLIDE 7: ARCHITECTURE DIAGRAM (Split Column) -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "System Architecture Blueprint", "06. TECHNICAL TOPOLOGY")

create_card(slide, 0.8, 1.8, 5.6, 1.15, 
            "Frontend (Vercel Host)", "Next.js App Router loaded locally on client viewports. Caches assets for offline capability.")

create_card(slide, 0.8, 3.05, 5.6, 1.15, 
            "REST API Server (Render Host)", "Node/Express endpoint handling geofence coordinates checks and messaging gateways.")

create_card(slide, 0.8, 4.3, 5.6, 1.15, 
            "NoSQL Database (Google Firestore)", "Stores student profiles, chat logs, stamps, and coordinates as document schemas.")

create_card(slide, 0.8, 5.55, 5.6, 1.15, 
            "File Host (Cloudinary CDN)", "Compresses and stores student scan uploads globally for mobile load speeds.")

add_placeholder_box(slide, 6.9, 1.8, 5.6, 4.9, "[ ARCHITECTURE DIAGRAM: FRONTEND / API / STORAGE / CLOUD ]")

# ----------------- SLIDE 8: DATA/USER FLOW DIAGRAM (Split Column) -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Data & User Flow Diagram", "07. PIPELINE FLOW")

create_card(slide, 0.8, 1.8, 5.6, 1.5, 
            "Location Scanning Pipeline", 
            "Camera Upload -> API uploads to Cloudinary -> API checks location databases -> Firestore logs stamp -> Client unlocks Map Tile.",
            title_color=COLOR_ACCENT_ORANGE)

create_card(slide, 0.8, 3.5, 5.6, 1.5, 
            "Explore Streak Calculation", 
            "Client logs scan -> Server checks lastScanDate -> if equal to 1 day, increments streak and awards +25 XP streak bonus.",
            title_color=COLOR_ACCENT_ORANGE)

create_card(slide, 0.8, 5.2, 5.6, 1.5, 
            "Databricks Genie Ingestion", 
            "Firestore records -> Syncs to GCS Bucket via Firebase Extension -> Delta Lake mount -> Databricks Genie query execution.",
            title_color=COLOR_ACCENT_ORANGE)

add_placeholder_box(slide, 6.9, 1.8, 5.6, 4.9, "[ DIAGRAM: DATA SEQUENCE TRANSACTION FLOW ]")

# ----------------- SLIDE 9: DATA MODEL USAGE (3 Column Layout) -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Data Model Usage", "08. DATABASE DESIGN")

col_width = 3.65
x_col1 = 0.8
x_col2 = 4.85
x_col3 = 8.9
height_col = 4.9

# Column 1
create_card(slide, x_col1, 1.8, col_width, height_col, 
            "Student Collections", 
            "•  users collection\nPath: `/users/{userId}`\nFields: name, username, avatarUrl, level, xp, interests, skills, exploreStreak, lastScanDate\n\n•  passport_stamps collection\nPath: `/users/{userId}/stamps/{stampId}`\nFields: locationId, unlockedAt, isFirstDocumenter, locationName",
            title_color=COLOR_TEXT_MAIN)

# Column 2
create_card(slide, x_col2, 1.8, col_width, height_col, 
            "Campus Collections", 
            "•  locations collection\nPath: `/locations/{locationId}`\nFields: name, building, floor, coordinates, rarity, facilities, tips\n\n•  map_tiles collection\nPath: `/map_tiles/{tileId}`\nFields: name, coordinates, iconType, associatedLocationId",
            title_color=COLOR_TEXT_MAIN)

# Column 3
create_card(slide, x_col3, 1.8, col_width, height_col, 
            "Social & Community", 
            "•  messages collection\nPath: `/chats/{chatId}/messages/{messageId}`\nFields: senderId, receiverId, text, imageUrl, sharedLocationId, timestamp\n\n•  posts collection\nPath: `/posts/{postId}`\nFields: authorName, title, content, tag, likes, timestamp",
            title_color=COLOR_TEXT_MAIN)

# ----------------- SLIDE 10: AI FEATURES AND GROUNDING (3 Column Layout) -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "AI Features & Grounding Model", "09. ARTIFICIAL INTELLIGENCE")

# 3 vertical columns representing our distinct AI components
create_card(slide, x_col1, 1.8, col_width, height_col, 
            "Databricks Genie Agent", 
            "Provides administrators with conversational SQL data analysis. Genie maps natural English to mounted database schemas, automatically translating queries like 'Which landmarks had the most scans?' into charts without coding queries.",
            title_color=COLOR_ACCENT_PURPLE)

create_card(slide, x_col2, 1.8, col_width, height_col, 
            "Gemini Grounded Chatbot", 
            "Powers the Genie Chatbot in the mobile client. It is grounded by system instructions seeded with campus guides, print rates, office locations, and schedules, answering student requests accurately.",
            title_color=COLOR_ACCENT_PURPLE)

create_card(slide, x_col3, 1.8, col_width, height_col, 
            "Profile Matchmaker", 
            "Uses mathematical sets calculations directly on Firestore collections. Computes overlap ratios across classmate interests, skills, and hobbies to matches students with high-score partners.",
            title_color=COLOR_ACCENT_PURPLE)

# ----------------- SLIDE 11: VALUE PROPOSITION AND BENEFITS (3 Column Layout) -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Value Proposition & Benefits", "10. STRATEGIC VALUE")

create_card(slide, x_col1, 1.8, col_width, height_col, 
            "For Students", 
            "•  Onboarding Speed: Freshmen discover campus facilities (printer spots, hardware labs) in minutes.\n\n•  Peer Matchmaking: Connects peers directly based on code skills or hobbies for hackathon sprints.\n\n•  Active Play: Gamified achievements make daily walks fun.",
            title_color=COLOR_TEXT_MAIN)

create_card(slide, x_col2, 1.8, col_width, height_col, 
            "For Admins", 
            "•  Natural Queries: Admins generate charts dynamically using Databricks Genie.\n\n•  Analytics: Reveals real-time spatial heatmaps showing popular spots.\n\n•  Resource Planning: Helps coordinates schedules based on facility traffic checks.",
            title_color=COLOR_TEXT_MAIN)

create_card(slide, x_col3, 1.8, col_width, height_col, 
            "For the Institution", 
            "•  Student Retention: Unifies student communication, reducing early isolation.\n\n•  Platform Consolidation: Integrates geolocations, chats, and forums into one app.\n\n•  Event Engagement: Drives participation in scavenger hunts and tech club events.",
            title_color=COLOR_TEXT_MAIN)

# ----------------- SLIDE 12: TECH STACK & PRODUCTION PLATFORMS (Grid List) -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Technology Stack & Free Tier Plan", "11. CLOUD STACK")

tech_cards = [
    ("Next.js & React (Vercel Host)", "Deploys static assets. Hobby free tier gives 100 GB monthly bandwidth and free SSL domain connection."),
    ("Google Firestore (Firebase Free Tier)", "Serverless NoSQL database. Free tier includes 50,000 reads, 20,000 writes daily, and 1 GiB storage."),
    ("Node.js & Express (Render Host)", "Runs backend geofence controls, explore streaks calculations, and real-time message streams."),
    ("Cloudinary (Media Hosting)", "Compresses and stores student scans. Free tier includes 25,000 monthly image transformations or 25 GB storage."),
    ("Databricks Genie Agent (CE)", "Databricks Community Edition processes analytics tables mounted from Firestore data lakes.")
]

y_tech = 1.8
for idx, (title, desc) in enumerate(tech_cards):
    create_card(slide, 0.8, y_tech, 11.7, 0.9, title, desc, title_color=COLOR_ACCENT_ORANGE)
    y_tech += 1.05

# ----------------- SLIDE 13: EXECUTION PLAN - 12 HOURS (Timeline Horizontal Flow) -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Execution Plan - 12 Hours", "12. ROADMAP")

timeline = [
    ("H 1 - 3: Setup & Database", "Config Next.js app, style theme files, set up Firebase Firestore document collections, and initialize auth APIs."),
    ("H 4 - 6: Backend & Upload", "Code backend Express APIs for photo scans, geofencing validations, Cloudinary storage uploads, and streak comparators."),
    ("H 7 - 9: Client Integration", "Bind React Context state logic inside views. Integrate Maps coordinates, messaging streams, and Hub profile matches."),
    ("H 10 - 12: Analytics & Host", "Sync Firestore to GCS, mount Delta tables in Databricks Genie, run compile builds, and deploy onto Vercel and Render.")
]

x_time = 0.8
width_time = 2.75
gap = 0.2
for idx, (title, desc) in enumerate(timeline):
    create_card(slide, x_time, 1.8, width_time, 4.9, title, desc, title_color=COLOR_ACCENT_PURPLE)
    x_time += width_time + gap

# ----------------- SLIDE 14: RISKS AND FALLBACKS (3 Column Layout) -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Risks and Fallbacks", "13. MITIGATION")

create_card(slide, x_col1, 1.8, col_width, height_col, 
            "Risk 1: Firestore Daily Read limits", 
            "•  Impact: App crashes once reads surpass the daily 50k limit under high student usage.\n\n•  Mitigation: Implement aggressive client-side caching via browser LocalStorage. Only fetch database collections when new chat pins or scans are updated, keeping reads minimal.",
            border_color=COLOR_ACCENT_ORANGE, title_color=COLOR_TEXT_MAIN)

create_card(slide, x_col2, 1.8, col_width, height_col, 
            "Risk 2: Render Free Tier Cold Starts", 
            "•  Impact: Rendering scanner triggers lags 30-50 seconds when user server restarts.\n\n•  Mitigation: Send a periodic uptime heartbeat query from the client application background to keep Render's free instance active during peak campus hours.",
            border_color=COLOR_ACCENT_ORANGE, title_color=COLOR_TEXT_MAIN)

create_card(slide, x_col3, 1.8, col_width, height_col, 
            "Risk 3: Internet Outage in Basements", 
            "•  Impact: Camera scans fail in basement labs due to dead cellular signals.\n\n•  Mitigation: Buffer scans in a local offline queue in LocalStorage, syncing coordinates and images once internet connections are restored.",
            border_color=COLOR_ACCENT_ORANGE, title_color=COLOR_TEXT_MAIN)

# Save presentation
prs.save("BMSCE_Campus_Social_Platform_Polished.pptx")
print("Polished Presentation generated successfully!")
