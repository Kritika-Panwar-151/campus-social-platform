import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
prs.slide_height = Inches(7.5)

# Style Palette Constants
COLOR_BG = RGBColor(248, 247, 255)         # Ermine White (#f8f7ff)
COLOR_TEXT_MAIN = RGBColor(51, 49, 54)     # Stiletto Charcoal (#333136)
COLOR_TEXT_SUB = RGBColor(99, 93, 115)     # Luscious Purple (#635d73)
COLOR_ACCENT_ORANGE = RGBColor(242, 143, 95) # Candied Yam (#f28f5f)
COLOR_ACCENT_PURPLE = RGBColor(162, 124, 248) # Purple Illusionist (#a27cf8)
COLOR_WHITE = RGBColor(255, 255, 255)

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header(slide, title_text, category_text="BMSCE CAMPUS SOCIAL PLATFORM"):
    # Add category tag
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = category_text.upper()
    p.font.name = "Arial"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_ORANGE
    
    # Add title
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = title_text
    p2.font.name = "Arial"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_MAIN

def add_placeholder_box(slide, x, y, width, height, label_text):
    # Add outline rectangle for screenshot / diagram placeholder
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 238, 250)
    shape.line.color.rgb = COLOR_ACCENT_PURPLE
    shape.line.width = Pt(1.5)
    
    # Text inside placeholder
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label_text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_SUB

# ----------------- SLIDE 1: TITLE SLIDE -----------------
slide_layout = prs.slide_layouts[6] # Blank
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)

# Title box
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.0))
tf = txBox.text_frame
tf.word_wrap = True

p_category = tf.paragraphs[0]
p_category.text = "NEXT-GENERATION STUDENT ENGAGEMENT"
p_category.font.name = "Arial"
p_category.font.size = Pt(11)
p_category.font.bold = True
p_category.font.color.rgb = COLOR_ACCENT_ORANGE
p_category.space_after = Pt(12)

p_title = tf.add_paragraph()
p_title.text = "BMSCE Campus Social & Intelligence Platform"
p_title.font.name = "Arial"
p_title.font.size = Pt(44)
p_title.font.bold = True
p_title.font.color.rgb = COLOR_TEXT_MAIN
p_title.space_after = Pt(8)

p_sub = tf.add_paragraph()
p_sub.text = "A Gamified, AI-Powered Offline-First Campus Ecosystem"
p_sub.font.name = "Arial"
p_sub.font.size = Pt(18)
p_sub.font.color.rgb = COLOR_TEXT_SUB

# ----------------- SLIDE 2: THE PROBLEM: A DISCONNECTED CAMPUS -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "The Problem: A Disconnected Campus", "CHALLENGES IN STUDENT LIFE")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True

points = [
    ("Freshman Isolation & Friction", "New students face a steep learning curve adapting to campus geography, building nomenclature, and academic clubs, leading to early social friction and isolation."),
    ("Fragmented Communication Channels", "Campus updates, hackathons, and workshops are scattered across email lists, chat groups, and physical notice boards, causing students to miss important events."),
    ("Lack of Target Peer Matching", "Finding classmates with matching technical skills (e.g., React, Python) or mutual hobbies (e.g., Chess, Sketching) for team sprints is highly unstructured and inefficient."),
    ("Static Navigation Tools", "Traditional campus layout maps are static PDFs that do not encourage exploration, leaving students unaware of key resources like hardware labs, printing kiosks, and quiet study zones.")
]

for idx, (headline, body) in enumerate(points):
    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
    p.text = f"•  {headline}: "
    p.font.name = "Arial"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.name = "Arial"
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_TEXT_SUB
    p2.space_after = Pt(8)

# ----------------- SLIDE 3: SCOPE: WHAT WE ARE BUILDING -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Scope: What We Are Building", "SOLUTION BOUNDARIES")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True

scope_points = [
    ("Gamified Campus Exploration Scanner", "A mobile-first web scanner allowing students to scan landmarks, verify locations, and unlock coordinate cards under distinct Rarity Tiers."),
    ("Fog-of-War Interactive Campus Map", "An interactive map starts greyed out and opens up tile-by-tile as students scan and visit physical blocks, gamifying campus layouts."),
    ("Campus Passport & Badge System", "A collection checkbook documenting stamps for scanned zones. Completing building sets unlocks badges and XP rank levels."),
    ("Skill-Based Social Matchmaking", "Direct student matching displaying overlapping hobbies and skills percentages, providing connections and sharing verified pins."),
    ("Databricks Genie AI Insights Desk", "An analytics panel translating conversational questions from admins into structured charts showing campus traffic patterns.")
]

for idx, (title, desc) in enumerate(scope_points):
    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
    p.text = f"✔  {title}"
    p.font.name = "Arial"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_PURPLE
    p.space_before = Pt(6)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = "Arial"
    p2.font.size = Pt(12.5)
    p2.font.color.rgb = COLOR_TEXT_MAIN
    p2.space_after = Pt(6)

# ----------------- SLIDE 4: BMSCE PLATFORM ECOSYSTEM -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "BMSCE Platform Ecosystem", "INTEGRATED STUDENT HUB")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True

eco_points = [
    ("Mobile Front-End Client", "A responsive app offering five interactive tabs: Explorer (map/scanner), Campus (matches), Genie (AI chatbot), Inbox (messaging), and Profile."),
    ("Serverless Backend", "Cloud-native middleware routing scan photo compression, streak dates, and real-time message sync."),
    ("Analytics Data Ingestion", "Firestore data is synced directly to a Delta Lakehouse environment, where administrators perform zero-code analytics using natural language."),
    ("Student Community Board", "A peer-to-peer forum filtered by categories like #Hackathons and #Workshops, helping students broadcast updates.")
]

for idx, (title, desc) in enumerate(eco_points):
    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_ORANGE
    p.space_before = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = "Arial"
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_TEXT_SUB
    p2.space_after = Pt(8)

add_placeholder_box(slide, 7.3, 1.8, 5.2, 4.8, "[ ECOSYSTEM ARCHITECTURE DIAGRAM / PREVIEW ]")

# ----------------- SLIDE 5: USER JOURNEY: EXPLORING THE CAMPUS -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "User Journey: Exploring the Campus", "GAMIFIED USER EXPERIENCE")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True

steps = [
    ("1. Check Locked Zones", "User inspects the Fog Map and finds the 'Engineering block C' sector coordinates greyed out."),
    ("2. Scan Location", "User visits the building, opens the Scanner, and snaps a photograph of the Robotics Lab."),
    ("3. Stamp Passport & Unlocks", "Firestore stores the scan, stamps their Passport with an official ink seal, and clears the fog to show a CPU icon."),
    ("4. Level Up Celebration", "The scan earns them +120 XP (First to Document) triggering the rank promotion overlay to level up.")
]

for idx, (title, desc) in enumerate(steps):
    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(6)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = "Arial"
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = COLOR_TEXT_SUB
    p2.space_after = Pt(6)

add_placeholder_box(slide, 7.3, 1.8, 5.2, 4.8, "[ SCREENSHOT: EXPLORER SCAN / FOG MAP / PASSPORT ]")

# ----------------- SLIDE 6: USER JOURNEY: FINDING YOUR TRIBE -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "User Journey: Finding Your Tribe", "SOCIAL MATCHMAKING EXPERIENCE")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True

social_steps = [
    ("1. Input Interest Tags", "User edits profile, selecting skills (React, Figma) and hobbies (Chess, Gaming)."),
    ("2. View Smart Match Scores", "The Matchmaker algorithm compares profiles and displays peers sorted by matching ratios (e.g. 85% overlap)."),
    ("3. Connect & Start Chatting", "User taps Connect to immediately trigger a peer notification and open a DM thread."),
    ("4. Pin Spot Coordinates", "User shares the coordinates of a newly unlocked 'Hidden Gem' study room in the chat, allowing the peer to view it instantly.")
]

for idx, (title, desc) in enumerate(social_steps):
    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(6)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = "Arial"
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = COLOR_TEXT_SUB
    p2.space_after = Pt(6)

add_placeholder_box(slide, 7.3, 1.8, 5.2, 4.8, "[ SCREENSHOT: HUB MATCHES / PRIVATE CHAT / DMs ]")

# ----------------- SLIDE 7: ARCHITECTURE DIAGRAM -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "System Architecture", "COMPONENTS & CLOUD PLACEMENT")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True

arch_desc = [
    ("Frontend (Vercel)", "Next.js pages cache assets locally. Connects securely via firebase-auth client integrations."),
    ("Backend API (Render/GCP)", "Runs Express controllers checking scan geofences, streak intervals, and message dispatches."),
    ("NoSQL Database (Firestore)", "Generous free-tier document collections storing user info, locations, and chat databases."),
    ("Cloud Storage (Cloudinary)", "Handles instant image compression, resizing, and global CDN hosting for scan uploads.")
]

for idx, (title, desc) in enumerate(arch_desc):
    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
    p.text = f"■  {title}"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_PURPLE
    p.space_before = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = "Arial"
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_TEXT_MAIN
    p2.space_after = Pt(8)

add_placeholder_box(slide, 7.3, 1.8, 5.2, 4.8, "[ DIAGRAM: SYSTEM ARCHITECTURE ]")

# ----------------- SLIDE 8: DATA/USER FLOW DIAGRAM -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Data & User Flow Diagram", "INFORMATION PIPELINES")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True

flow_desc = [
    ("Location Scan Flow", "Client takes photo -> API uploads it to Cloudinary -> API checks first-ever scan count -> Firestore locks/unlocks the passport stamp."),
    ("Streak heartbeats Flow", "Firestore compares timestamp dates -> if exactly 24h since last scan, increments streak and triggers +25 XP bonus."),
    ("Databricks Genie Ingestion Flow", "Firestore collections -> Syncs to GCS Bucket via Firebase Extension -> Mounted as Delta Tables in Databricks -> Genie Agent parses query requests.")
]

for idx, (title, desc) in enumerate(flow_desc):
    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_ORANGE
    p.space_before = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = "Arial"
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = COLOR_TEXT_SUB
    p2.space_after = Pt(8)

add_placeholder_box(slide, 7.3, 1.8, 5.2, 4.8, "[ DIAGRAM: DATA SEQUENCE FLOW ]")

# ----------------- SLIDE 9: DATA MODEL USAGE -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Data Model Usage", "FIRESTORE NOSQL COLLECTIONS")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True

models = [
    ("1. users Collection", "Paths: `/users/{userId}`. Fields: `name`, `username`, `avatarUrl`, `level`, `xp`, `interests`, `skills`, `exploreStreak`, `lastScanDate`."),
    ("2. locations Collection", "Paths: `/locations/{locationId}`. Fields: `name`, `building`, `floor`, `coordinates`, `rarity`, `facilities`, `tips`."),
    ("3. passport_stamps Collection", "Paths: `/users/{userId}/stamps/{stampId}`. Fields: `locationId`, `unlockedAt`, `isFirstDocumenter`, `locationName`."),
    ("4. messages Collection", "Paths: `/chats/{chatId}/messages/{messageId}`. Fields: `senderId`, `receiverId`, `text`, `imageUrl`, `sharedLocationId`, `timestamp`."),
    ("5. discussion_posts Collection", "Paths: `/posts/{postId}`. Fields: `authorName`, `title`, `content`, `tag`, `likes`, `timestamp`.")
]

for idx, (title, desc) in enumerate(models):
    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(6)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = "Arial"
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_TEXT_SUB
    p2.space_after = Pt(6)

# ----------------- SLIDE 10: AI FEATURES AND GROUNDING -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "AI Features and Grounding", "INTELLIGENT CAMPUS AGENTS")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True

ai_features = [
    ("Databricks Genie Analytics Agent", "Grounds admin queries using schema context files. Translates questions like 'Which buildings are trending?' into optimized Delta Lakehouse SQL commands, plotting tables dynamically without coding errors."),
    ("Google Gemini Chatbot (Genie Tab)", "Integrates Gemini SDK within the app. Grounded by system prompt injection containing college guidelines, printing rates, club contacts, and directories to answer student questions accurately."),
    ("Automated Peer Matchmaker", "Heuristic filtering computes overlap ratios in user profiles to suggest partners, mapping matching percentages on student cards for project sprints.")
]

for idx, (title, desc) in enumerate(ai_features):
    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
    p.text = f"★  {title}"
    p.font.name = "Arial"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_ORANGE
    p.space_before = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = "Arial"
    p2.font.size = Pt(12.5)
    p2.font.color.rgb = COLOR_TEXT_MAIN
    p2.space_after = Pt(8)

# ----------------- SLIDE 11: VALUE PROPOSITION AND BENEFITS -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Value Proposition and Benefits", "IMPACT AND VALUE CREATION")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True

benefits = [
    ("For Students: Gamified Integration & Peer Matching", "Speeds up campus onboarding for freshers. Provides an intuitive way to discover facilities (printer spots, quiet rooms) and find relevant study/hackathon teams."),
    ("For Admins: Decentered Analytics & Flow Insights", "Saves time generating charts. Using Databricks Genie, admins query user scan coordinates and traffic logs through plain text to identify bottleneck hotspots on campus."),
    ("For the Institution: Increased Engagement", "Increases student event attendance and builds a collaborative, supportive campus community by unifying chats, forum boards, and scavenger events.")
]

for idx, (title, desc) in enumerate(benefits):
    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
    p.text = f"●  {title}"
    p.font.name = "Arial"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = "Arial"
    p2.font.size = Pt(12.5)
    p2.font.color.rgb = COLOR_TEXT_SUB
    p2.space_after = Pt(8)

# ----------------- SLIDE 12: TECH STACK -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Technology Stack", "FREE-TIER PRODUCTION SYSTEM")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True

stack_details = [
    ("Next.js 14 & React 18 (Vercel Hosting)", "Free tier handles asset loading, page transitions, and animations with 100 GB bandwidth monthly."),
    ("Google Cloud Firestore (Firebase Free Storage)", "Generous database limits: 50k reads, 20k writes, 20k deletes daily, and 1 GiB data storage."),
    ("Node.js & Express (Render Hosting)", "Hosts middle layer API for location geofencing, explore streaks, and messaging triggers."),
    ("Cloudinary Media Storage", "Saves student camera scans. Free tier includes 25,000 monthly image transformations or 25 GB storage."),
    ("Databricks Genie Agent (Community Edition)", "Queries analytical tables mounted via GCS, translating plain English into Lakehouse SQL charts.")
]

for idx, (title, desc) in enumerate(stack_details):
    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
    p.text = f"✔  {title}"
    p.font.name = "Arial"
    p.font.size = Pt(14.5)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_PURPLE
    p.space_before = Pt(6)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = "Arial"
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_TEXT_MAIN
    p2.space_after = Pt(6)

# ----------------- SLIDE 13: EXECUTION PLAN - 12 HOURS -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Execution Plan - 12 Hours", "RAPID DEPLOYMENT ROADMAP")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True

execution_plan = [
    ("Hours 1 - 3: Environment Setup & Auth initialization", "Initialize Next.js app, configure Tailwind styles, set up Firebase Console, activate Firebase Auth, and seed the Firestore database structure."),
    ("Hours 4 - 6: Core Gamification API Development", "Code backend Express logic for camera uploads, Cloudinary storage setup, passport stamps logging, and streak comparisons."),
    ("Hours 7 - 9: Client Interface Integration", "Wire up React Context hook inside frontend views. Bind Scanner maps, messaging streams, and Hub profile recommendations."),
    ("Hours 10 - 12: Analytics desk config, polishing & Vercel deployment", "Sync Firestore to GCS, mount Delta tables in Databricks Genie, verify code compile builds, and deploy onto Vercel and Render free hosts.")
]

for idx, (title, desc) in enumerate(execution_plan):
    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(14.5)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = "Arial"
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_TEXT_SUB
    p2.space_after = Pt(8)

# ----------------- SLIDE 14: RISKS AND FALLBACKS -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Risks and Fallbacks", "MITIGATION STRATEGIES")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True

risks = [
    ("Risk 1: Firestore Daily Read Limit Exceeded", "Mitigation: Implement aggressive client-side caching using browser LocalStorage. Only sync modified documents (messages, stamp items) to Firestore when new scans or chat updates occur."),
    ("Risk 2: Render API Server Cold Starts", "Mitigation: Set up a periodic uptime heartbeat query from the client application or a cron scheduler to prevent Render's free instance from sleeping, ensuring immediate response times."),
    ("Risk 3: Internet Outage in Basements", "Mitigation: Queue scan data locally in a LocalStorage buffer while offline, and automatically upload and process coordinates once connection is restored.")
]

for idx, (title, desc) in enumerate(risks):
    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
    p.text = f"⚠  {title}"
    p.font.name = "Arial"
    p.font.size = Pt(14.5)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_ORANGE
    p.space_before = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = "Arial"
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_TEXT_MAIN
    p2.space_after = Pt(8)

# Save presentation
prs.save("BMSCE_Campus_Social_Platform.pptx")
print("Presentation generated successfully!")
